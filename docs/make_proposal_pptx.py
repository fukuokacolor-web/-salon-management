# -*- coding: utf-8 -*-
"""
サロン顧客管理システム 提案資料 (.pptx) 生成スクリプト
出力: SALON_PROPOSAL.pptx (16:9, 20スライド)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ========== カラー定義 ==========
ROSE      = RGBColor(0xD4, 0x68, 0x8A)  # メインカラー
DARK_ROSE = RGBColor(0x7B, 0x4A, 0x5C)  # アクセント濃
GOLD      = RGBColor(0xD4, 0xA5, 0x74)  # ゴールド
BG_PINK   = RGBColor(0xFF, 0xF5, 0xF7)  # 極薄ピンク
LIGHT_PINK= RGBColor(0xFF, 0xE4, 0xEC)  # 薄ピンク
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY= RGBColor(0xEE, 0xEE, 0xEE)
GREEN     = RGBColor(0x5C, 0xB8, 0x5C)
ORANGE    = RGBColor(0xF0, 0xA0, 0x40)
RED       = RGBColor(0xD9, 0x53, 0x4F)
LINE_GREEN= RGBColor(0x06, 0xC7, 0x55)

FONT_JP = 'メイリオ'
FONT_EN = 'Arial'

# ========== プレゼン初期化 ==========
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]

shape_count = 0
table_count = 0

# ========== ヘルパー ==========
def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    global shape_count
    shape_count += 1
    return bg

def add_top_bottom_line(slide):
    """上下の細いピンクライン"""
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Emu(45720))  # 上ライン
    top.line.fill.background()
    top.fill.solid(); top.fill.fore_color.rgb = ROSE
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Emu(45720), SW, Emu(45720))
    bot.line.fill.background()
    bot.fill.solid(); bot.fill.fore_color.rgb = ROSE
    global shape_count
    shape_count += 2

def add_text(slide, x, y, w, h, text, *,
             font=FONT_JP, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """テキストボックスを追加。text に \\n を入れると段落分割。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        # 日本語フォント明示
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rpr, qn('a:ea'))
        ea.set('typeface', FONT_JP)
    global shape_count
    shape_count += 1
    return tb

def add_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=1.0,
             shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    global shape_count
    shape_count += 1
    return s

def add_rounded(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=1.0):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)

def add_title_bar(slide, title_text, num):
    """共通ヘッダ: 左に番号丸、タイトル、上下ライン"""
    add_top_bottom_line(slide)
    # 番号丸
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(0.3),
                                    Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = ROSE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"{num:02d}"
    r.font.name = FONT_EN; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WHITE
    global shape_count
    shape_count += 1
    # タイトル
    add_text(slide, Inches(1.15), Inches(0.28), Inches(11), Inches(0.7),
             title_text, size=26, bold=True, color=DARK_ROSE,
             anchor=MSO_ANCHOR.MIDDLE)
    # 下線（タイトル下）
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4),
                                  Inches(1.05), Inches(12.5), Emu(18000))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    shape_count += 1

def add_footer(slide, num, total=20):
    """右下にページ番号、左下に商品名"""
    add_text(slide, Inches(0.3), Inches(7.05), Inches(6), Inches(0.3),
             "サロン顧客管理システム ご提案資料",
             size=9, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.6), Inches(0.3),
             f"{num} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ========================================================
# Slide 1 — 表紙
# ========================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, BG_PINK)
# 装飾帯
add_rect(s, 0, Inches(0), SW, Inches(0.25), fill=ROSE)
add_rect(s, 0, Inches(7.25), SW, Inches(0.25), fill=ROSE)
# 中央大きな丸
deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.67), Inches(1.2),
                          Inches(2), Inches(2))
deco.fill.solid(); deco.fill.fore_color.rgb = ROSE
deco.line.fill.background()
shape_count += 1
add_text(s, Inches(5.67), Inches(1.2), Inches(2), Inches(2), "🌸",
         size=72, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         color=WHITE)
# タイトル
add_text(s, Inches(1), Inches(3.5), Inches(11.33), Inches(1.2),
         "LINEで予約が完結する",
         size=36, bold=True, color=DARK_ROSE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.3), Inches(11.33), Inches(1.2),
         "サロン顧客管理システム",
         size=44, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
# サブタイトル
add_text(s, Inches(1), Inches(5.5), Inches(11.33), Inches(0.6),
         "オーナー様向け ご提案資料",
         size=20, color=DARK_ROSE, align=PP_ALIGN.CENTER)
# バージョン
add_text(s, Inches(1), Inches(6.4), Inches(11.33), Inches(0.4),
         "Version 1.0  |  2026.XX",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ========================================================
# Slide 2 — こんなお悩みありませんか？
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "こんなお悩みありませんか？", 2)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "一つでも当てはまれば、このシステムで解決できます。",
         size=16, color=DARK_ROSE)
worries = [
    "予約の電話・LINEメッセージ対応に追われている",
    "無断キャンセルや常連の予約忘れが多い",
    "顧客ごとのコース残回数の管理が大変",
    "エクセル管理でデータが散らばっている",
    "スマホで業務確認したいが既存システムは重い",
]
top = 2.0
for i, w in enumerate(worries):
    y = Inches(top + i * 0.95)
    card = add_rounded(s, Inches(1.0), y, Inches(11.3), Inches(0.8),
                       fill=BG_PINK, line=ROSE, line_w=1.0)
    # チェックボックス
    add_rect(s, Inches(1.3), y + Inches(0.18), Inches(0.44), Inches(0.44),
             fill=WHITE, line=ROSE, line_w=1.5)
    add_text(s, Inches(1.3), y + Inches(0.13), Inches(0.44), Inches(0.5),
             "□", size=22, bold=True, color=ROSE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.95), y, Inches(10.2), Inches(0.8),
             w, size=18, bold=True, color=BLACK,
             anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 2)

# ========================================================
# Slide 3 — このシステムが解決します
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "このシステムが解決します", 3)
# メインコピー
add_rounded(s, Inches(1.5), Inches(1.5), Inches(10.33), Inches(1.2),
            fill=ROSE)
add_text(s, Inches(1.5), Inches(1.5), Inches(10.33), Inches(1.2),
         "予約・顧客・ポイント・売上。すべてを一つに。",
         size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
# 3本柱
cols = [
    ("📱", "LINE予約",     "24時間いつでも受付\nアプリ不要"),
    ("💻", "管理画面",     "スマホ・PC対応\nどこでも確認"),
    ("⚡", "自動化",       "リマインド・バックアップ\n手間なし"),
]
for i, (icon, title, desc) in enumerate(cols):
    x = Inches(0.7 + i * 4.15)
    y = Inches(3.2)
    card = add_rounded(s, x, y, Inches(4.0), Inches(3.3),
                       fill=BG_PINK, line=ROSE, line_w=2.0)
    add_text(s, x, y + Inches(0.3), Inches(4.0), Inches(1.0),
             icon, size=54, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.4), Inches(4.0), Inches(0.6),
             title, size=22, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.1), Inches(4.0), Inches(1.0),
             desc, size=14, color=BLACK, align=PP_ALIGN.CENTER)
add_footer(s, 3)

# ========================================================
# Slide 4 — 3つの特徴
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "3つの特徴", 4)
features = [
    ("01", "LINEから24時間予約（アプリ不要）",
     "お客様は普段お使いのLINEで、予約・変更・キャンセル・残回数確認がすべて完結。\n新しいアプリをダウンロードする必要はありません。"),
    ("02", "スマホ対応の管理画面",
     "オーナー様はスマホ・タブレット・PCから業務を確認・操作。\nお客様のご来店中でも、外出先でもスムーズに対応できます。"),
    ("03", "自動リマインド・自動バックアップ",
     "前日に自動でお客様へリマインドLINE。週次で自動バックアップ。\n「うっかり忘れ」も「データ消失」も、仕組みで防ぎます。"),
]
for i, (no, title, body) in enumerate(features):
    y = Inches(1.35 + i * 1.85)
    add_rounded(s, Inches(0.5), y, Inches(12.3), Inches(1.65),
                fill=BG_PINK, line=ROSE, line_w=1.5)
    # 番号
    num_box = add_rect(s, Inches(0.5), y, Inches(1.5), Inches(1.65),
                       fill=ROSE)
    add_text(s, Inches(0.5), y, Inches(1.5), Inches(1.65),
             no, font=FONT_EN, size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.2), y + Inches(0.2), Inches(10.3), Inches(0.5),
             title, size=18, bold=True, color=DARK_ROSE)
    add_text(s, Inches(2.2), y + Inches(0.75), Inches(10.3), Inches(0.85),
             body, size=12, color=BLACK, line_spacing=1.3)
add_footer(s, 4)

# ========================================================
# Slide 5 — 顧客体験①：LINEで予約完結
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "顧客体験① LINEで予約完結", 5)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "お客様はLINEアプリひとつで、24時間いつでも予約できます。",
         size=15, color=DARK_ROSE)

# LINE画面モック（3画面）
def line_phone(x, y, w, h, title_txt, body_fn):
    # フレーム
    frame = add_rounded(s, x, y, w, h, fill=LIGHT_GRAY, line=DARK_ROSE, line_w=2)
    # ヘッダ（LINE緑）
    hdr = add_rect(s, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2),
                   Inches(0.4), fill=LINE_GREEN)
    add_text(s, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.4),
             title_txt, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 画面エリア
    scr = add_rect(s, x + Inches(0.1), y + Inches(0.55),
                   w - Inches(0.2), h - Inches(0.65), fill=WHITE)
    body_fn(x + Inches(0.1), y + Inches(0.55), w - Inches(0.2), h - Inches(0.65))

def body1(bx, by, bw, bh):
    # リッチメニュー風（上段2ボタン＋下段フルワイド / 実装と同じ3ボタン構成）
    add_text(s, bx, by + Inches(0.1), bw, Inches(0.4),
             "リッチメニュー", size=9, color=GRAY,
             align=PP_ALIGN.CENTER)
    # 上段: 予約する / 予約確認
    top_btns = [("予約する", ROSE), ("予約確認", GOLD)]
    top_w = Inches(1.3); top_h = Inches(0.95)
    for i, (t, c) in enumerate(top_btns):
        bx2 = bx + Inches(0.15) + Inches(i * 1.4)
        by2 = by + Inches(1.45)
        add_rounded(s, bx2, by2, top_w, top_h, fill=c)
        add_text(s, bx2, by2, top_w, top_h, t,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 下段: お問い合わせ（フルワイド）
    full_w = Inches(2.7); full_h = Inches(0.95)
    add_rounded(s, bx + Inches(0.15), by + Inches(2.5),
                full_w, full_h, fill=DARK_ROSE)
    add_text(s, bx + Inches(0.15), by + Inches(2.5),
             full_w, full_h, "お問い合わせ",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def body2(bx, by, bw, bh):
    # カレンダーグリッド
    add_text(s, bx, by + Inches(0.1), bw, Inches(0.3),
             "2026年4月", size=10, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    days = ['日','月','火','水','木','金','土']
    cw = bw / 7
    # 曜日ヘッダ
    for i, d in enumerate(days):
        add_text(s, bx + cw * i, by + Inches(0.45), cw, Inches(0.2),
                 d, size=8, color=GRAY, align=PP_ALIGN.CENTER)
    # 日付グリッド（3週分）
    marks = [['○','×','○','○','-','○','○'],
             ['○','○','×','○','○','-','○'],
             ['-','○','○','×','○','○','○']]
    nums = [['1','2','3','4','5','6','7'],
            ['8','9','10','11','12','13','14'],
            ['15','16','17','18','19','20','21']]
    for r in range(3):
        for c in range(7):
            mx = bx + cw * c
            my = by + Inches(0.7) + Inches(r * 0.5)
            m = marks[r][c]
            col = GREEN if m == '○' else (RED if m == '×' else GRAY)
            add_text(s, mx, my, cw, Inches(0.22),
                     nums[r][c], size=7, color=BLACK, align=PP_ALIGN.CENTER)
            add_text(s, mx, my + Inches(0.18), cw, Inches(0.22),
                     m, size=9, bold=True, color=col, align=PP_ALIGN.CENTER)

def body3(bx, by, bw, bh):
    add_text(s, bx, by + Inches(0.1), bw, Inches(0.3),
             "4月8日（水）", size=10, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    times = ["10:00", "11:00", "13:00 ✓", "14:00", "15:00", "16:00"]
    for i, t in enumerate(times):
        sel = "✓" in t
        y2 = by + Inches(0.5 + i * 0.38)
        add_rounded(s, bx + Inches(0.15), y2, bw - Inches(0.3), Inches(0.3),
                    fill=ROSE if sel else WHITE,
                    line=ROSE, line_w=0.75)
        add_text(s, bx + Inches(0.15), y2, bw - Inches(0.3), Inches(0.3),
                 t, size=10, bold=sel,
                 color=WHITE if sel else BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

line_phone(Inches(0.6), Inches(2.0), Inches(3.5), Inches(4.5),
           "① メニューを開く", body1)
line_phone(Inches(4.9), Inches(2.0), Inches(3.5), Inches(4.5),
           "② 日付を選ぶ", body2)
line_phone(Inches(9.2), Inches(2.0), Inches(3.5), Inches(4.5),
           "③ 時間を選んで確定", body3)

# 矢印
for ax in [Inches(4.25), Inches(8.55)]:
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, Inches(4.0),
                            Inches(0.5), Inches(0.5))
    ar.fill.solid(); ar.fill.fore_color.rgb = ROSE
    ar.line.fill.background()
    shape_count += 1

# 吹き出し
callout = add_rounded(s, Inches(3.5), Inches(6.6), Inches(6.3), Inches(0.5),
                     fill=GOLD)
add_text(s, Inches(3.5), Inches(6.6), Inches(6.3), Inches(0.5),
         "💡 アプリをインストールする必要ナシ！",
         size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 5)

# ========================================================
# Slide 6 — 顧客体験②：予約変更・キャンセル・残回数確認
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "顧客体験② 予約変更・キャンセル・残回数確認", 6)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "お客様の「ちょっとした用件」も、LINEで完結します。",
         size=15, color=DARK_ROSE)
items = [
    ("🔄", "予約変更",
     "日時の変更がLINEから可能。\n予約確認メッセージの\n「変更する」ボタン1タップ。"),
    ("❌", "予約キャンセル",
     "急な用事でもLINEで即キャンセル。\nキャンセル料ポリシーも\n自動で事前表示。"),
    ("📊", "コース残回数",
     "「残回数」ボタンで現在の残数を\nいつでも確認。\n次回予約の目安にも。"),
]
for i, (icon, t, d) in enumerate(items):
    x = Inches(0.5 + i * 4.3)
    y = Inches(2.0)
    add_rounded(s, x, y, Inches(4.0), Inches(3.8),
                fill=WHITE, line=ROSE, line_w=2)
    # アイコン円
    ic = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.25),
                            y + Inches(0.3), Inches(1.5), Inches(1.5))
    ic.fill.solid(); ic.fill.fore_color.rgb = BG_PINK
    ic.line.color.rgb = ROSE; ic.line.width = Pt(1.5)
    shape_count += 1
    add_text(s, x + Inches(1.25), y + Inches(0.3), Inches(1.5), Inches(1.5),
             icon, size=48, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y + Inches(2.0), Inches(4.0), Inches(0.5),
             t, size=20, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(2.6), Inches(3.6), Inches(1.2),
             d, size=12, color=BLACK, align=PP_ALIGN.CENTER)

# 結論バー
add_rounded(s, Inches(1.5), Inches(6.3), Inches(10.33), Inches(0.6),
            fill=ROSE)
add_text(s, Inches(1.5), Inches(6.3), Inches(10.33), Inches(0.6),
         "💬 わざわざ電話する必要がない ＝ お客様の満足度UP",
         size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 6)

# ========================================================
# Slide 7 — オーナー体験①：ダッシュボード
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "オーナー体験① ダッシュボード", 7)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "本日必要な情報が、ひと目でわかる画面設計。",
         size=15, color=DARK_ROSE)

# ダッシュボード画面モック
add_rounded(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(4.9),
            fill=BG_PINK, line=ROSE, line_w=2)
# ヘッダ
add_rect(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.5),
         fill=ROSE)
add_text(s, Inches(0.9), Inches(2.05), Inches(10), Inches(0.5),
         "🌸 〇〇サロン  ダッシュボード",
         size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(10.5), Inches(2.05), Inches(2), Inches(0.5),
         "2026/04/21 ▼",
         size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

# 4カード（統計）— 実装のダッシュボード表示項目に準拠
stats = [
    ("本日の予約", "5件",     ROSE),
    ("本日の売上", "¥38,500", GOLD),
    ("今月の売上", "¥482,000", DARK_ROSE),
    ("コース残少", "3名",     ORANGE),
]
for i, (t, v, c) in enumerate(stats):
    x = Inches(0.9 + i * 2.95)
    y = Inches(2.75)
    add_rounded(s, x, y, Inches(2.75), Inches(1.3),
                fill=WHITE, line=c, line_w=2)
    add_text(s, x, y + Inches(0.1), Inches(2.75), Inches(0.35),
             t, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(2.75), Inches(0.7),
             v, size=24, bold=True, color=c,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 本日の予約リスト
add_rounded(s, Inches(0.9), Inches(4.2), Inches(7.5), Inches(2.5),
            fill=WHITE, line=LIGHT_PINK, line_w=1)
add_rect(s, Inches(0.9), Inches(4.2), Inches(7.5), Inches(0.35),
         fill=LIGHT_PINK)
add_text(s, Inches(1.05), Inches(4.2), Inches(7), Inches(0.35),
         "📅 本日の予約",
         size=12, bold=True, color=DARK_ROSE, anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("10:00", "田中 美咲様",   "カット＋カラー"),
    ("11:30", "佐藤 健一様",   "フェイシャル"),
    ("13:00", "鈴木 陽子様",   "コース (残3回)"),
    ("15:00", "高橋 由美様",   "カット"),
]
for i, (tm, nm, mn) in enumerate(rows):
    yy = Inches(4.6 + i * 0.5)
    add_text(s, Inches(1.1), yy, Inches(1.2), Inches(0.4),
             tm, size=11, bold=True, color=ROSE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.4), yy, Inches(2.5), Inches(0.4),
             nm, size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.0), yy, Inches(3.2), Inches(0.4),
             mn, size=11, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)

# 右: バックアップ＆コース残回数アラート
add_rounded(s, Inches(8.6), Inches(4.2), Inches(4.0), Inches(1.1),
            fill=WHITE, line=GREEN, line_w=1.5)
add_text(s, Inches(8.7), Inches(4.25), Inches(3.8), Inches(0.4),
         "💾 最新バックアップ ✓", size=11, bold=True, color=GREEN)
add_text(s, Inches(8.7), Inches(4.65), Inches(3.8), Inches(0.6),
         "2026/04/19 02:00\n（毎週日曜 午前2時に自動実行）",
         size=10, color=BLACK)

add_rounded(s, Inches(8.6), Inches(5.4), Inches(4.0), Inches(1.3),
            fill=WHITE, line=ORANGE, line_w=1.5)
add_text(s, Inches(8.7), Inches(5.45), Inches(3.8), Inches(0.4),
         "⚠ コース残回数アラート", size=11, bold=True, color=ORANGE)
add_text(s, Inches(8.7), Inches(5.85), Inches(3.8), Inches(0.8),
         "残り4回以下: 3名\n→ 次回来店時にご案内を",
         size=10, color=BLACK, line_spacing=1.3)
add_footer(s, 7)

# ========================================================
# Slide 8 — オーナー体験②：顧客管理・コース残回数
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "オーナー体験② 顧客管理・コース残回数", 8)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "残回数が少なくなったお客様を、自動でピックアップ。",
         size=15, color=DARK_ROSE)

# 顧客一覧テーブル
add_rounded(s, Inches(0.5), Inches(2.0), Inches(8.3), Inches(4.8),
            fill=WHITE, line=ROSE, line_w=1.5)
add_rect(s, Inches(0.5), Inches(2.0), Inches(8.3), Inches(0.5),
         fill=ROSE)
add_text(s, Inches(0.7), Inches(2.0), Inches(8), Inches(0.5),
         "👥 顧客一覧",
         size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# 列ヘッダ
hdr_y = Inches(2.6)
headers = [("顧客ID", 1.1), ("お名前", 2.5), ("コース", 2.2), ("残回数", 1.3), ("来店", 1.0)]
x_pos = Inches(0.6)
for h, ww in headers:
    add_text(s, x_pos, hdr_y, Inches(ww), Inches(0.35),
             h, size=10, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    x_pos += Inches(ww)

# データ行
data = [
    ("C001", "田中 美咲", "フェイシャル10回", "8回", GREEN, "5回"),
    ("C002", "佐藤 健一", "ボディケア5回",    "4回", ORANGE, "2回"),
    ("C003", "鈴木 陽子", "ヘアトリートメント10回", "1回", RED, "8回"),
    ("C004", "高橋 由美", "カラー5回",        "3回", ORANGE, "3回"),
    ("C005", "山田 花子", "コース終了",        "-",  GRAY,   "12回"),
    ("C006", "伊藤 明子", "フェイシャル20回", "15回", GREEN, "6回"),
    ("C007", "渡辺 梨花", "ボディケア10回",   "2回", RED,    "9回"),
]
for i, (cid, nm, cs, rn, col, vs) in enumerate(data):
    yy = Inches(3.0 + i * 0.5)
    if i % 2 == 0:
        add_rect(s, Inches(0.55), yy, Inches(8.2), Inches(0.45), fill=BG_PINK)
    xp = Inches(0.6)
    for v, ww, c, sz in [(cid, 1.1, DARK_ROSE, 10), (nm, 2.5, BLACK, 11),
                          (cs, 2.2, BLACK, 10)]:
        add_text(s, xp, yy, Inches(ww), Inches(0.45),
                 v, size=sz, color=c, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        xp += Inches(ww)
    # 残回数バッジ
    badge = add_rounded(s, xp + Inches(0.2), yy + Inches(0.08),
                        Inches(0.9), Inches(0.3), fill=col)
    add_text(s, xp + Inches(0.2), yy + Inches(0.08), Inches(0.9), Inches(0.3),
             rn, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    xp += Inches(1.3)
    add_text(s, xp, yy, Inches(1.0), Inches(0.45),
             vs, size=10, color=BLACK, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)

# 右サイドアラート
add_rounded(s, Inches(9.0), Inches(2.0), Inches(3.8), Inches(2.3),
            fill=RED, line=None)
add_text(s, Inches(9.0), Inches(2.1), Inches(3.8), Inches(0.5),
         "⚠️ 要フォロー", size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(9.1), Inches(2.6), Inches(3.6), Inches(1.6),
         "残り4回以下の\nお客様が3名います。\n\nコース追加のご案内を\nお忘れなく！",
         size=12, color=WHITE, align=PP_ALIGN.CENTER, line_spacing=1.3)

# バッジ凡例
add_rounded(s, Inches(9.0), Inches(4.5), Inches(3.8), Inches(2.3),
            fill=WHITE, line=ROSE, line_w=1.5)
add_text(s, Inches(9.0), Inches(4.6), Inches(3.8), Inches(0.5),
         "🏷️ バッジ色の見方", size=14, bold=True, color=DARK_ROSE,
         align=PP_ALIGN.CENTER)
legends = [("5回以上", GREEN, "余裕あり"),
           ("2〜4回",  ORANGE, "そろそろ案内"),
           ("1回以下", RED,    "要フォロー")]
for i, (lv, c, ls) in enumerate(legends):
    yy = Inches(5.15 + i * 0.5)
    add_rounded(s, Inches(9.2), yy, Inches(0.9), Inches(0.35), fill=c)
    add_text(s, Inches(9.2), yy, Inches(0.9), Inches(0.35),
             lv, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.3), yy, Inches(2.4), Inches(0.35),
             ls, size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 8)

# ========================================================
# Slide 9 — オーナー体験③：売上レポート
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "オーナー体験③ 売上レポート", 9)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "月別・日別の売上が、グラフで一目瞭然。売上明細・顧客一覧・予約台帳はCSVでダウンロード可能（Excel対応）。",
         size=14, color=DARK_ROSE)

# 左：棒グラフ
add_rounded(s, Inches(0.5), Inches(2.0), Inches(7.5), Inches(4.8),
            fill=WHITE, line=ROSE, line_w=1.5)
add_text(s, Inches(0.7), Inches(2.1), Inches(7), Inches(0.4),
         "📊 月別売上（過去12ヶ月）",
         size=13, bold=True, color=DARK_ROSE)
# 棒グラフ描画
months = ["5","6","7","8","9","10","11","12","1","2","3","4"]
values = [0.45, 0.52, 0.48, 0.55, 0.6, 0.65, 0.7, 0.82, 0.58, 0.62, 0.75, 0.88]
graph_x = Inches(0.9); graph_y_base = Inches(6.0); graph_h_max = Inches(2.8)
bar_w = Inches(0.45); gap = Inches(0.12)
for i, (m, v) in enumerate(zip(months, values)):
    bh = Emu(int(int(graph_h_max) * v))
    bx = graph_x + (bar_w + gap) * i
    by = graph_y_base - bh
    col = ROSE if v >= 0.8 else (GOLD if v >= 0.6 else LIGHT_PINK)
    bar = add_rect(s, bx, by, bar_w, bh, fill=col)
    add_text(s, bx - Inches(0.02), Inches(6.05), bar_w + Inches(0.04),
             Inches(0.3), m, size=9, color=GRAY, align=PP_ALIGN.CENTER)
# Y軸
add_text(s, Inches(0.55), Inches(2.7), Inches(0.35), Inches(0.3),
         "¥", size=10, color=GRAY)
add_text(s, Inches(0.6), Inches(6.35), Inches(6.8), Inches(0.3),
         "月", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# 右：CSV風
add_rounded(s, Inches(8.2), Inches(2.0), Inches(4.6), Inches(4.8),
            fill=WHITE, line=GOLD, line_w=1.5)
add_text(s, Inches(8.4), Inches(2.1), Inches(4.3), Inches(0.4),
         "📁 CSVダウンロード（売上／顧客／予約）",
         size=12, bold=True, color=DARK_ROSE)
# テーブル風
csv_headers = ["日付", "顧客", "金額"]
csv_rows = [
    ("04/20", "田中様", "¥8,500"),
    ("04/20", "佐藤様", "¥12,000"),
    ("04/19", "鈴木様", "¥6,000"),
    ("04/19", "高橋様", "¥15,000"),
    ("04/18", "山田様", "¥8,500"),
    ("04/18", "伊藤様", "¥22,000"),
    ("04/17", "渡辺様", "¥8,500"),
]
# ヘッダ行
add_rect(s, Inches(8.4), Inches(2.6), Inches(4.2), Inches(0.35),
         fill=ROSE)
for i, h in enumerate(csv_headers):
    add_text(s, Inches(8.4 + i * 1.4), Inches(2.6), Inches(1.4), Inches(0.35),
             h, size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
for r, row in enumerate(csv_rows):
    yy = Inches(2.95 + r * 0.32)
    if r % 2 == 0:
        add_rect(s, Inches(8.4), yy, Inches(4.2), Inches(0.3), fill=BG_PINK)
    for i, v in enumerate(row):
        add_text(s, Inches(8.4 + i * 1.4), yy, Inches(1.4), Inches(0.3),
                 v, size=9, color=BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# DLボタン
add_rounded(s, Inches(8.8), Inches(5.9), Inches(3.4), Inches(0.6),
            fill=GOLD)
add_text(s, Inches(8.8), Inches(5.9), Inches(3.4), Inches(0.6),
         "⬇ CSVダウンロード",
         size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 9)

# ========================================================
# Slide 10 — スマホ完全対応
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "スマホ完全対応", 10)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "PCでもスマホでも、画面サイズに合わせて自動で最適化。",
         size=15, color=DARK_ROSE)

# 左：PC画面
add_rounded(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.2),
            fill=BLACK)
add_rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(3.3),
         fill=WHITE)
add_rect(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(0.3), fill=ROSE)
# テーブル
tb_headers = ["日時", "顧客名", "メニュー", "金額"]
for i, h in enumerate(tb_headers):
    add_text(s, Inches(0.9 + i * 1.35), Inches(2.75),
             Inches(1.35), Inches(0.3),
             h, size=9, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
pc_rows = [("4/21 10:00","田中様","カット","¥5,500"),
           ("4/21 11:30","佐藤様","カラー","¥12,000"),
           ("4/21 13:00","鈴木様","フェイシャル","¥8,000"),
           ("4/21 15:00","高橋様","コース","-")]
for r, row in enumerate(pc_rows):
    yy = Inches(3.1 + r * 0.35)
    for i, v in enumerate(row):
        add_text(s, Inches(0.9 + i * 1.35), yy, Inches(1.35), Inches(0.35),
                 v, size=8, color=BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# スタンド
add_rect(s, Inches(2.9), Inches(6.3), Inches(1.4), Inches(0.15), fill=BLACK)
add_rect(s, Inches(2.3), Inches(6.45), Inches(2.6), Inches(0.1), fill=BLACK)
add_text(s, Inches(0.6), Inches(6.6), Inches(6.0), Inches(0.3),
         "PC：広々と一覧表示",
         size=12, bold=True, color=DARK_ROSE, align=PP_ALIGN.CENTER)

# 右：スマホ画面
add_rounded(s, Inches(8.9), Inches(2.0), Inches(2.8), Inches(4.6),
            fill=BLACK)
add_rect(s, Inches(9.0), Inches(2.25), Inches(2.6), Inches(4.1),
         fill=WHITE)
add_rect(s, Inches(9.0), Inches(2.25), Inches(2.6), Inches(0.3), fill=ROSE)
add_text(s, Inches(9.0), Inches(2.25), Inches(2.6), Inches(0.3),
         "予約一覧", size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# カード型
for r in range(4):
    yy = Inches(2.65 + r * 0.88)
    add_rounded(s, Inches(9.1), yy, Inches(2.4), Inches(0.8),
                fill=BG_PINK, line=ROSE, line_w=0.75)
    add_text(s, Inches(9.2), yy + Inches(0.05), Inches(2.3), Inches(0.3),
             f"4/21 {10 + r * 2}:00", size=9, bold=True, color=DARK_ROSE)
    nms = ["田中様","佐藤様","鈴木様","高橋様"]
    mns = ["カット ¥5,500","カラー ¥12,000","フェイシャル ¥8,000","コース"]
    add_text(s, Inches(9.2), yy + Inches(0.3), Inches(2.3), Inches(0.25),
             nms[r], size=9, color=BLACK)
    add_text(s, Inches(9.2), yy + Inches(0.52), Inches(2.3), Inches(0.25),
             mns[r], size=8, color=GRAY)
add_text(s, Inches(8.9), Inches(6.75), Inches(2.8), Inches(0.3),
         "スマホ：カード型に自動変換",
         size=11, bold=True, color=DARK_ROSE, align=PP_ALIGN.CENTER)

# 中央矢印＋メッセージ
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.0),
                        Inches(3.8), Inches(1.7), Inches(0.8))
ar.fill.solid(); ar.fill.fore_color.rgb = GOLD
ar.line.fill.background()
shape_count += 1
add_text(s, Inches(7.0), Inches(3.8), Inches(1.7), Inches(0.8),
         "自動\n切替", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 10)

# ========================================================
# Slide 11 — 自動リマインド＆自動バックアップ
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "自動リマインド & 自動バックアップ", 11)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "「うっかり」と「もしもの時」から、サロンを守る仕組み。",
         size=15, color=DARK_ROSE)

# 左：リマインド
add_rounded(s, Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.5),
            fill=BG_PINK, line=ROSE, line_w=2)
add_text(s, Inches(0.7), Inches(2.15), Inches(5.8), Inches(0.5),
         "⏰ 前日リマインド", size=20, bold=True, color=ROSE)
add_text(s, Inches(0.7), Inches(2.7), Inches(5.8), Inches(0.4),
         "お客様のLINEに自動でメッセージ",
         size=12, color=DARK_ROSE)
# LINE吹き出し
add_rounded(s, Inches(0.9), Inches(3.2), Inches(5.2), Inches(1.8),
            fill=WHITE, line=GRAY, line_w=0.75)
add_text(s, Inches(1.1), Inches(3.3), Inches(5.0), Inches(1.6),
         "〇〇サロンです🌸\n\n明日 4/22（金）14:00 から\nカット＋カラーのご予約を\nお待ちしております。\n\n変更・キャンセルはメニューから。",
         size=11, color=BLACK, line_spacing=1.3)
# 効果
add_rounded(s, Inches(0.7), Inches(5.15), Inches(5.7), Inches(1.15),
            fill=GREEN)
add_text(s, Inches(0.7), Inches(5.2), Inches(5.7), Inches(0.4),
         "効果", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(5.55), Inches(5.7), Inches(0.6),
         "無断キャンセル激減 → 機会損失を防止",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 右：バックアップ
add_rounded(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.5),
            fill=BG_PINK, line=GOLD, line_w=2)
add_text(s, Inches(7.1), Inches(2.15), Inches(5.7), Inches(0.5),
         "💾 週次バックアップ", size=20, bold=True, color=GOLD)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.7), Inches(0.4),
         "毎週日曜 午前2時、自動でデータを保管",
         size=12, color=DARK_ROSE)
# スケジュールイメージ
sched = [("日","2:00", True), ("月","", False), ("火","", False),
         ("水","", False), ("木","", False), ("金","", False), ("土","", False)]
for i, (d, tm, active) in enumerate(sched):
    x = Inches(7.15 + i * 0.8)
    c = ROSE if active else LIGHT_GRAY
    add_rounded(s, x, Inches(3.2), Inches(0.7), Inches(0.9),
                fill=c)
    add_text(s, x, Inches(3.25), Inches(0.7), Inches(0.3),
             d, size=10, bold=True,
             color=WHITE if active else GRAY,
             align=PP_ALIGN.CENTER)
    if tm:
        add_text(s, x, Inches(3.55), Inches(0.7), Inches(0.35),
                 tm, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
add_text(s, Inches(7.1), Inches(4.25), Inches(5.8), Inches(0.4),
         "→ ダッシュボードに「最新バックアップ✓」表示",
         size=11, color=BLACK)
# バックアップ成功表示
add_rounded(s, Inches(7.3), Inches(4.7), Inches(5.3), Inches(0.45),
            fill=WHITE, line=GREEN, line_w=1.5)
add_text(s, Inches(7.3), Inches(4.7), Inches(5.3), Inches(0.45),
         "✓ 2026/04/19 02:00 バックアップ完了",
         size=12, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 効果
add_rounded(s, Inches(7.1), Inches(5.3), Inches(5.7), Inches(1.0),
            fill=GOLD)
add_text(s, Inches(7.1), Inches(5.3), Inches(5.7), Inches(1.0),
         "💡 データ消失の心配なし\n過去4週間分を自動保管",
         size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

add_footer(s, 11)

# ========================================================
# Slide 12 — セキュリティ対策
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "セキュリティ対策", 12)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "大切なお客様情報を、しっかり守る4つの仕組み。",
         size=15, color=DARK_ROSE)

sec_items = [
    ("🔐", "パスワードはハッシュ化保存",
     "SHA-256で不可逆変換して保存。\n万一データが見られても\nパスワードは復元不可能です。"),
    ("🔒", "通信はすべて暗号化",
     "LINE・管理画面とサーバー間の\n通信はHTTPSで暗号化。\n盗聴の心配はありません。"),
    ("👤", "共有PCでも安心",
     "ログイン情報の保存を選べる方式。\nOFFならブラウザを閉じた時点で\n自動ログアウトします。"),
    ("📝", "操作履歴を自動記録",
     "管理者ログイン履歴と\nLINEメッセージ送受信履歴を\n自動記録。後から確認できます。"),
]
for i, (icon, title, desc) in enumerate(sec_items):
    r = i // 2; c = i % 2
    x = Inches(0.5 + c * 6.4)
    y = Inches(2.0 + r * 2.4)
    add_rounded(s, x, y, Inches(6.2), Inches(2.2),
                fill=WHITE, line=ROSE, line_w=1.5)
    # アイコン
    ic = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2),
                            y + Inches(0.3), Inches(1.4), Inches(1.4))
    ic.fill.solid(); ic.fill.fore_color.rgb = ROSE
    ic.line.fill.background()
    shape_count += 1
    add_text(s, x + Inches(0.2), y + Inches(0.3), Inches(1.4), Inches(1.4),
             icon, size=40, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.8), y + Inches(0.3), Inches(4.2), Inches(0.5),
             title, size=16, bold=True, color=DARK_ROSE)
    add_text(s, x + Inches(1.8), y + Inches(0.85), Inches(4.3), Inches(1.2),
             desc, size=11, color=BLACK, line_spacing=1.3)
add_footer(s, 12)

# ========================================================
# Slide 13 — 他システムとの違い（テーブル）
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "他システムとの違い", 13)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "コスト・柔軟性・データ所有権。すべてサロン側のメリットを追求。",
         size=15, color=DARK_ROSE)

# python-pptx Table
rows_data = [
    ["項目", "当システム", "大手予約サイト"],
    ["月額固定費", "要相談（変動なし）", "高額 ＋ 手数料"],
    ["LINE連携", "標準搭載", "オプション"],
    ["カスタマイズ\n（店舗名・文言・料金）", "自由に変更可", "制限あり"],
    ["データ所有", "サロン側", "事業者側"],
    ["予約手数料", "なし", "1件ごとに課金"],
]
tbl_x = Inches(1.0); tbl_y = Inches(2.0)
tbl_w = Inches(11.3); tbl_h = Inches(4.5)
tbl_shape = s.shapes.add_table(len(rows_data), 3, tbl_x, tbl_y, tbl_w, tbl_h)
tbl = tbl_shape.table
table_count += 1
shape_count += 1

# 列幅
tbl.columns[0].width = Inches(3.3)
tbl.columns[1].width = Inches(4.0)
tbl.columns[2].width = Inches(4.0)

for ri, row in enumerate(rows_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = val
        run.font.name = FONT_JP
        if ri == 0:
            run.font.size = Pt(16); run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = ROSE
        else:
            run.font.size = Pt(14)
            if ci == 0:
                run.font.bold = True
                run.font.color.rgb = DARK_ROSE
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_PINK
            elif ci == 1:
                run.font.bold = True
                run.font.color.rgb = ROSE
                cell.fill.solid(); cell.fill.fore_color.rgb = BG_PINK
            else:
                run.font.color.rgb = GRAY
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
        # 日本語EAフォント
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rpr, qn('a:ea'))
        ea.set('typeface', FONT_JP)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.35),
         "※ 上記は一般的な比較です。具体的な金額・条件は個別にご相談ください。",
         size=10, color=GRAY)
add_footer(s, 13)

# ========================================================
# Slide 14 — 導入の流れ（4週間）
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "導入の流れ（4週間）", 14)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "お申込みから運用開始まで、およそ4週間で完了します。",
         size=15, color=DARK_ROSE)

weeks = [
    ("Week 1", "ヒアリング & 環境準備",
     "サロン様の業務フローをお伺いし、\nGoogleアカウント・LINE公式アカウント\nの準備をサポート。"),
    ("Week 2", "LINE公式アカウント連携",
     "予約メニュー・リッチメニューの設定、\nサロン専用の挨拶メッセージや\nキャンセルポリシーを登録。"),
    ("Week 3", "初期データ移行・動作確認",
     "既存顧客データの取り込み、\nメニュー・コース料金の登録、\nテスト予約での動作確認。"),
    ("Week 4", "運用開始 & 操作研修",
     "オンライン研修（管理画面・\nよくある操作）を実施。\n本格運用スタート後もサポート継続。"),
]
# 横並び4ステップ
for i, (wk, ttl, dsc) in enumerate(weeks):
    x = Inches(0.45 + i * 3.15)
    y = Inches(2.0)
    # 番号円
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1),
                                y, Inches(0.9), Inches(0.9))
    circle.fill.solid(); circle.fill.fore_color.rgb = ROSE
    circle.line.color.rgb = GOLD; circle.line.width = Pt(2.5)
    shape_count += 1
    add_text(s, x + Inches(1.1), y, Inches(0.9), Inches(0.9),
             str(i + 1), font=FONT_EN, size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # カード
    add_rounded(s, x, y + Inches(1.1), Inches(3.0), Inches(3.6),
                fill=WHITE, line=ROSE, line_w=1.5)
    add_text(s, x, y + Inches(1.2), Inches(3.0), Inches(0.4),
             wk, font=FONT_EN, size=14, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.65), Inches(3.0), Inches(0.6),
             ttl, size=14, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(2.4), Inches(2.7),
             Inches(2.2), dsc, size=11, color=BLACK,
             line_spacing=1.3, align=PP_ALIGN.CENTER)
    # 矢印（最終以外）
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                x + Inches(3.0), y + Inches(2.5),
                                Inches(0.2), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD
        ar.line.fill.background()
        shape_count += 1

add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.35),
         "※ サロン様のご都合やデータ量により、期間は前後する場合があります。",
         size=10, color=GRAY)
add_footer(s, 14)

# ========================================================
# Slide 15 — 料金プラン
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "料金プラン", 15)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "サロン様の規模・ご要望に合わせて、個別にお見積もりいたします。",
         size=15, color=DARK_ROSE)

# 3枚カード
plans = [
    ("💰", "初期費用",         "要相談",
     "• 環境構築・初期設定\n• LINE公式アカウント連携\n• 既存データ移行\n• オンライン研修1回"),
    ("📅", "月額保守費",       "要相談",
     "• システム稼働監視\n• データバックアップ\n• LINE・メールサポート\n• 軽微な修正対応"),
    ("⭐", "オプション",       "要相談",
     "• カスタマイズ開発\n• 追加研修\n• 繁忙期サポート強化\n• 複数店舗対応"),
]
for i, (icon, ttl, pr, dsc) in enumerate(plans):
    x = Inches(0.5 + i * 4.3)
    y = Inches(2.0)
    mid = (i == 1)
    fill_c = ROSE if mid else WHITE
    text_c = WHITE if mid else DARK_ROSE
    desc_c = WHITE if mid else BLACK
    add_rounded(s, x, y, Inches(4.0), Inches(4.5),
                fill=fill_c, line=ROSE, line_w=2)
    add_text(s, x, y + Inches(0.3), Inches(4.0), Inches(0.8),
             icon, size=44, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.3), Inches(4.0), Inches(0.5),
             ttl, size=20, bold=True, color=text_c,
             align=PP_ALIGN.CENTER)
    # 区切り線
    ln = add_rect(s, x + Inches(0.5), y + Inches(1.85),
                  Inches(3.0), Emu(9000),
                  fill=GOLD if not mid else WHITE)
    # 価格
    add_text(s, x, y + Inches(2.0), Inches(4.0), Inches(0.8),
             pr, size=32, bold=True,
             color=WHITE if mid else ROSE,
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.3), y + Inches(2.9), Inches(3.5), Inches(1.4),
             dsc, size=11, color=desc_c, line_spacing=1.4)

# 注記
add_rounded(s, Inches(2), Inches(6.65), Inches(9.3), Inches(0.4),
            fill=BG_PINK)
add_text(s, Inches(2), Inches(6.65), Inches(9.3), Inches(0.4),
         "※ 個別見積もりいたします。お気軽にご相談ください。",
         size=12, bold=True, color=DARK_ROSE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 15)

# ========================================================
# Slide 16 — サポート体制
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "サポート体制", 16)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "導入後も安心の4段階サポート。困った時はいつでもご相談ください。",
         size=15, color=DARK_ROSE)

supports = [
    ("🎯", "導入時",
     "設定代行 ＋ オンライン研修", "初期設定を代行し、管理画面の操作方法を\nオンライン（画面共有）で丁寧にご説明。"),
    ("💬", "運用開始後",
     "LINE・メールサポート", "ご不明点はLINEまたはメールで\n24時間受付。営業日に順次回答します。"),
    ("📅", "月次",
     "運用状況の確認", "月1回、予約件数・売上・稼働状況を\nレポートとしてご報告。改善提案も。"),
    ("🚨", "緊急時",
     "障害対応", "システム停止や重大な不具合発生時は\n優先対応。復旧までの情報共有も密に。"),
]
for i, (icon, when, what, desc) in enumerate(supports):
    r = i // 2; c = i % 2
    x = Inches(0.5 + c * 6.4)
    y = Inches(2.0 + r * 2.4)
    add_rounded(s, x, y, Inches(6.2), Inches(2.2),
                fill=BG_PINK, line=ROSE, line_w=1.5)
    # 左アイコン
    ic = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25),
                            y + Inches(0.45), Inches(1.3), Inches(1.3))
    ic.fill.solid(); ic.fill.fore_color.rgb = GOLD
    ic.line.fill.background()
    shape_count += 1
    add_text(s, x + Inches(0.25), y + Inches(0.45), Inches(1.3), Inches(1.3),
             icon, size=36, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # タイトル
    add_text(s, x + Inches(1.75), y + Inches(0.2), Inches(4.3), Inches(0.4),
             when, size=12, bold=True, color=GOLD)
    add_text(s, x + Inches(1.75), y + Inches(0.55), Inches(4.3), Inches(0.5),
             what, size=16, bold=True, color=DARK_ROSE)
    add_text(s, x + Inches(1.75), y + Inches(1.1), Inches(4.4), Inches(1.0),
             desc, size=11, color=BLACK, line_spacing=1.3)
add_footer(s, 16)

# ========================================================
# Slide 17 — よくある質問
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "よくあるご質問", 17)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
         "導入前によくいただくご質問をまとめました。",
         size=14, color=DARK_ROSE)

qas = [
    ("操作が難しそうで不安です。",
     "ご心配いりません。導入時にマニュアルをお渡しし、オンライン研修で実際に操作しながらご説明します。運用開始後もいつでも質問OK。"),
    ("データはどこに保管されますか？",
     "Google社のクラウド（Googleドライブ）上に、オーナー様のアカウントで保管されます。データの所有権は完全にサロン様にあります。"),
    ("既存の顧客データはどうなりますか？",
     "現時点では売上明細・顧客一覧・予約台帳のCSV出力に対応しています。他システムからのデータ移行は個別にご相談のうえ対応いたします。"),
    ("途中で解約したくなったら？",
     "データはそのまま保持され、サロン様のGoogleアカウントに残ります。いつでも内容を確認・ダウンロードしていただけます。"),
    ("月の予約件数に制限はありますか？",
     "システム自体に件数上限はありません。LINE公式アカウントの無料枠を超える規模の場合は、有料プランへの移行をご相談いたします。"),
]
for i, (q, a) in enumerate(qas):
    y = Inches(1.85 + i * 1.02)
    # Q
    qbox = add_rounded(s, Inches(0.5), y, Inches(0.6), Inches(0.4), fill=ROSE)
    add_text(s, Inches(0.5), y, Inches(0.6), Inches(0.4),
             "Q", font=FONT_EN, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.2), y, Inches(11.5), Inches(0.4),
             q, size=14, bold=True, color=DARK_ROSE,
             anchor=MSO_ANCHOR.MIDDLE)
    # A
    abox = add_rounded(s, Inches(0.5), y + Inches(0.45),
                      Inches(0.6), Inches(0.4), fill=GOLD)
    add_text(s, Inches(0.5), y + Inches(0.45), Inches(0.6), Inches(0.4),
             "A", font=FONT_EN, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.2), y + Inches(0.45), Inches(11.5), Inches(0.45),
             a, size=11, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 17)

# ========================================================
# Slide 18 — お客様の声
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "お客様の声", 18)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "導入いただいたサロン様からのコメントをご紹介します。",
         size=15, color=DARK_ROSE)

voices = [
    ("🌸", "サロンA様（ヘアサロン）",
     "導入1ヶ月で、無断キャンセルが\n半分以下に減りました。\n前日リマインドの効果は絶大です。",
     "★★★★★"),
    ("✨", "サロンB様（エステ）",
     "お客様から『LINEで予約できて\n便利』と好評。新規のご予約も\n月30％増加しました。",
     "★★★★★"),
    ("💎", "サロンC様（ネイル）",
     "残回数が可視化されたことで、\nコース追加のご案内が\n自然にできるようになりました。",
     "★★★★★"),
]
for i, (icon, who, body, stars) in enumerate(voices):
    x = Inches(0.5 + i * 4.3)
    y = Inches(2.0)
    add_rounded(s, x, y, Inches(4.0), Inches(4.5),
                fill=WHITE, line=ROSE, line_w=2)
    # アイコン
    add_text(s, x, y + Inches(0.3), Inches(4.0), Inches(1.0),
             icon, size=54, align=PP_ALIGN.CENTER)
    # 星
    add_text(s, x, y + Inches(1.4), Inches(4.0), Inches(0.4),
             stars, size=16, color=GOLD, align=PP_ALIGN.CENTER)
    # サロン名
    add_text(s, x, y + Inches(1.85), Inches(4.0), Inches(0.4),
             who, size=13, bold=True, color=DARK_ROSE,
             align=PP_ALIGN.CENTER)
    # コメント
    add_text(s, x + Inches(0.3), y + Inches(2.35), Inches(3.4), Inches(2.0),
             "「" + body + "」", size=12, color=BLACK,
             align=PP_ALIGN.CENTER, line_spacing=1.4)

# 注釈
add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.35),
         "※ 上記はイメージです。実際の導入事例は、ご要望に応じて個別にご紹介いたします。",
         size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s, 18)

# ========================================================
# Slide 19 — 導入までの次のステップ
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "導入までの次のステップ", 19)
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "たった5ステップで、新しいサロン運営が始まります。",
         size=15, color=DARK_ROSE)

steps = [
    ("1", "本資料ご確認", "まずはこの資料をゆっくりご確認ください。"),
    ("2", "ヒアリング（30分無料）",
     "オンラインまたは訪問で、現状の課題をお伺いします。"),
    ("3", "お見積り提示",
     "サロン様に最適なプラン・金額をご提案。"),
    ("4", "ご契約・初期設定",
     "ご契約後、4週間で環境構築・データ移行を実施。"),
    ("5", "運用開始",
     "操作研修を経て、本格運用スタート！"),
]
# 縦並びタイムライン
tl_x = Inches(1.5)
# 縦線
add_rect(s, tl_x - Inches(0.01), Inches(2.1),
         Emu(18000), Inches(4.6), fill=GOLD)
for i, (no, ttl, dsc) in enumerate(steps):
    y = Inches(2.0 + i * 0.95)
    # 円
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, tl_x - Inches(0.4),
                                y, Inches(0.8), Inches(0.8))
    circle.fill.solid(); circle.fill.fore_color.rgb = ROSE
    circle.line.color.rgb = WHITE; circle.line.width = Pt(2)
    shape_count += 1
    add_text(s, tl_x - Inches(0.4), y, Inches(0.8), Inches(0.8),
             no, font=FONT_EN, size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # カード
    add_rounded(s, Inches(2.2), y, Inches(10.5), Inches(0.8),
                fill=BG_PINK, line=ROSE, line_w=1)
    add_text(s, Inches(2.4), y, Inches(3.5), Inches(0.8),
             ttl, size=15, bold=True, color=DARK_ROSE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.9), y, Inches(6.6), Inches(0.8),
             dsc, size=12, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 19)

# ========================================================
# Slide 20 — 連絡先
# ========================================================
s = prs.slides.add_slide(BLANK); add_bg(s, BG_PINK)
add_rect(s, 0, 0, SW, Inches(0.25), fill=ROSE)
add_rect(s, 0, Inches(7.25), SW, Inches(0.25), fill=ROSE)

# メインコピー
add_text(s, Inches(1), Inches(0.9), Inches(11.33), Inches(1.2),
         "まずはお気軽にご相談ください",
         size=36, bold=True, color=DARK_ROSE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(1.85), Inches(11.33), Inches(0.5),
         "Contact Us",
         font=FONT_EN, size=18, color=GOLD, align=PP_ALIGN.CENTER)

# 中央カード
add_rounded(s, Inches(2.5), Inches(2.7), Inches(8.3), Inches(3.8),
            fill=WHITE, line=ROSE, line_w=3)

# 担当者
add_text(s, Inches(2.5), Inches(3.0), Inches(8.3), Inches(0.4),
         "担当者", size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(3.35), Inches(8.3), Inches(0.5),
         "〇〇 〇〇",
         size=22, bold=True, color=DARK_ROSE, align=PP_ALIGN.CENTER)

# 区切り
add_rect(s, Inches(4), Inches(4.0), Inches(5.3), Emu(9000), fill=GOLD)

# メール
add_text(s, Inches(2.5), Inches(4.15), Inches(8.3), Inches(0.4),
         "📧 Email", size=12, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(4.5), Inches(8.3), Inches(0.4),
         "xxxxx@example.com",
         font=FONT_EN, size=18, color=BLACK, align=PP_ALIGN.CENTER)

# 電話
add_text(s, Inches(2.5), Inches(5.1), Inches(8.3), Inches(0.4),
         "📞 TEL", size=12, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(5.45), Inches(8.3), Inches(0.4),
         "000-0000-0000",
         font=FONT_EN, size=18, color=BLACK, align=PP_ALIGN.CENTER)

# 営業時間
add_text(s, Inches(2.5), Inches(6.05), Inches(8.3), Inches(0.4),
         "営業時間: 平日 10:00〜18:00",
         size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Thank you
add_text(s, Inches(1), Inches(6.7), Inches(11.33), Inches(0.5),
         "🌸 Thank you 🌸",
         font=FONT_EN, size=16, bold=True, color=ROSE, align=PP_ALIGN.CENTER)

# ========================================================
# 保存
# ========================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "SALON_PROPOSAL.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
print(f"Total shapes drawn (approx): {shape_count}")
print(f"Total tables: {table_count}")
